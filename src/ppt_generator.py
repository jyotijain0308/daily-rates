"""PPT Generator - Create country-specific product price list presentations."""
import logging
from datetime import datetime
from pathlib import Path
from typing import Optional

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config import (
    SLIDE_WIDTH, SLIDE_HEIGHT, COLORS, FONTS, FONT_SIZES, MARGINS,
    COMPANY_NAME, COMPANY_ADDRESS, COMPANY_WEBSITE, COMPANY_LOGO_IMAGE,
    UAE_LOGO_IMAGE, COUNTRY_LOGO_IMAGES, CURRENCY, RATE_DISPLAY_FORMAT
)
from product_data import ProductData
from product_image_service import ProductImageService

logger = logging.getLogger(__name__)


class PPTGenerator:
    """Generate PowerPoint presentations with product price lists."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.product_images = ProductImageService()

    def _resolve_asset(self, image_path: Optional[str]) -> Optional[Path]:
        if not image_path:
            return None

        path = Path(image_path)
        if path.is_absolute() and path.exists():
            return path

        candidates = [
            Path.cwd() / image_path,
            Path.cwd() / "src" / image_path,
            Path(__file__).resolve().parent / image_path,
            Path(__file__).resolve().parent.parent / image_path,
        ]
        for candidate in candidates:
            if candidate.exists():
                return candidate
        return None

    def _add_image_or_text(self, slide, image_path: Optional[str], fallback_text: str,
                           left, top, width, height, dark: bool = False,
                           fit_to_box: bool = False):
        image = self._resolve_asset(image_path)
        if image:
            try:
                if fit_to_box:
                    image_width, image_height = self._native_image_size(image)
                    scale = min(width / image_width, height / image_height)
                    display_width = int(image_width * scale)
                    display_height = int(image_height * scale)
                    picture = slide.shapes.add_picture(
                        str(image),
                        int(left + (width - display_width) / 2),
                        int(top + (height - display_height) / 2),
                        width=display_width,
                        height=display_height,
                    )
                else:
                    slide.shapes.add_picture(str(image), left, top, width=width)
                return
            except Exception as exc:
                logger.warning("Could not add image %s: %s", image, exc)

        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = fallback_text
        p.font.size = FONT_SIZES["small"]
        p.font.bold = True
        p.font.name = FONTS["body"]
        p.font.color.rgb = COLORS["light_text"] if dark else COLORS["primary"]
        p.alignment = PP_ALIGN.CENTER

    def _native_image_size(self, image_path: Path):
        with Image.open(image_path) as image:
            pixel_width, pixel_height = image.size
            dpi_x, dpi_y = image.info.get("dpi", (96, 96))

        dpi_x = dpi_x or 96
        dpi_y = dpi_y or 96
        return Inches(pixel_width / dpi_x), Inches(pixel_height / dpi_y)

    def _set_background(self, slide, color_key: str):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS[color_key]

    def _product_image_path(self, product_name: str) -> Optional[str]:
        return self.product_images.get_product_image_path(product_name)

    def add_country_title_slide(self, country_name: str, current_date: datetime = None):
        """Create country-specific title slide."""
        if current_date is None:
            current_date = datetime.now()

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, "background")

        self._add_image_or_text(
            slide, COMPANY_LOGO_IMAGE, COMPANY_NAME,
            Inches(0.5), Inches(0.45), Inches(9.0), Inches(0.95),
            fit_to_box=True
        )
        self._add_image_or_text(
            slide, COUNTRY_LOGO_IMAGES.get(country_name), country_name,
            Inches(0.5), Inches(1.95), Inches(9.0), Inches(2.2),
            fit_to_box=True
        )

        title_box = slide.shapes.add_textbox(
            MARGINS["left"], Inches(4.35),
            SLIDE_WIDTH - MARGINS["left"] - MARGINS["right"], Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = f"{country_name} Products Price List"
        p.font.size = FONT_SIZES["title"]
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.font.name = FONTS["title"]
        p.alignment = PP_ALIGN.CENTER

        date_box = slide.shapes.add_textbox(
            MARGINS["left"], Inches(5.85),
            SLIDE_WIDTH - MARGINS["left"] - MARGINS["right"], Inches(0.6)
        )
        date_frame = date_box.text_frame
        p = date_frame.paragraphs[0]
        p.text = current_date.strftime("%B %d, %Y")
        p.font.size = Pt(30)
        p.font.color.rgb = COLORS["text"]
        p.font.name = FONTS["body"]
        p.alignment = PP_ALIGN.CENTER

        logger.info("Added country title slide: %s", country_name)

    def add_title_slide(self, heading: str = "Daily Product Rates", current_date: datetime = None):
        """Backward-compatible title slide entry point."""
        self.add_country_title_slide(heading.replace(" Products Price List", ""), current_date)

    def add_product_slide(self, product: ProductData, slide_number: int = None):
        """Create one product item page."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, "background")

        self._add_image_or_text(
            slide, COMPANY_LOGO_IMAGE, COMPANY_NAME,
            Inches(0.35), Inches(0.25), Inches(1.25), Inches(0.6)
        )
        self._add_image_or_text(
            slide, COUNTRY_LOGO_IMAGES.get(product.country_of_origin), product.country_of_origin,
            Inches(8.55), Inches(0.25), Inches(1.0), Inches(0.65),
            fit_to_box=True
        )

        title_box = slide.shapes.add_textbox(
            MARGINS["left"], Inches(1.05),
            SLIDE_WIDTH - MARGINS["left"] - MARGINS["right"], Inches(1.15)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = f"{product.product_name} {product.weight_kg:g}kg {product.packing}"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.font.name = FONTS["title"]
        p.alignment = PP_ALIGN.CENTER

        self._add_image_or_text(
            slide, self._product_image_path(product.product_name), product.product_name,
            Inches(1.25), Inches(2.25), Inches(7.5), Inches(2.25),
            fit_to_box=True
        )

        price_band = slide.shapes.add_shape(
            1, Inches(1.35), Inches(5.15), Inches(7.3), Inches(1.0)
        )
        price_band.fill.solid()
        price_band.fill.fore_color.rgb = COLORS["primary"]
        price_band.line.color.rgb = COLORS["primary"]

        price_box = slide.shapes.add_textbox(
            Inches(1.35), Inches(5.33), Inches(7.3), Inches(0.55)
        )
        price_frame = price_box.text_frame
        p = price_frame.paragraphs[0]
        p.text = f"Price: {RATE_DISPLAY_FORMAT.format(product.price_aed)}"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = COLORS["light_text"]
        p.font.name = FONTS["title"]
        p.alignment = PP_ALIGN.CENTER

        logger.info("Added product slide: %s", product.product_name)

    def add_thank_you_slide(self, country_name: str, exchange_rate: Optional[float] = None,
                            currency_code: Optional[str] = None):
        """Create thank-you slide for a country PPT."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, "background")

        self._add_image_or_text(
            slide, COMPANY_LOGO_IMAGE, COMPANY_NAME,
            Inches(0.5), Inches(0.45), Inches(9.0), Inches(0.95),
            fit_to_box=True
        )

        self._add_image_or_text(
            slide, COUNTRY_LOGO_IMAGES.get(country_name), country_name,
            Inches(0.55), Inches(2.0), Inches(1.9), Inches(1.25),
            fit_to_box=True
        )
        self._add_image_or_text(
            slide, UAE_LOGO_IMAGE, "UAE",
            Inches(7.55), Inches(2.0), Inches(1.9), Inches(1.25),
            fit_to_box=True
        )

        rate_text = "Exchange rate not available"
        if exchange_rate and currency_code:
            rate_text = f"1 {CURRENCY} = {exchange_rate:.2f} {currency_code}"

        rate_background = slide.shapes.add_shape(
            1, Inches(3.15), Inches(2.25), Inches(3.7), Inches(0.74)
        )
        rate_background.fill.solid()
        rate_background.fill.fore_color.rgb = RGBColor(255, 244, 229)
        rate_background.line.color.rgb = COLORS["accent"]

        rate_box = slide.shapes.add_textbox(
            Inches(3.22), Inches(2.37), Inches(3.56), Inches(0.44)
        )
        rate_frame = rate_box.text_frame
        rate_frame.word_wrap = True
        rate_frame.margin_top = 0
        rate_frame.margin_bottom = 0
        rate_frame.margin_left = 0
        rate_frame.margin_right = 0
        p = rate_frame.paragraphs[0]
        p.text = rate_text
        p.font.size = FONT_SIZES["heading"]
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.font.name = FONTS["subtitle"]
        p.alignment = PP_ALIGN.CENTER

        company_box = slide.shapes.add_textbox(
            MARGINS["left"], Inches(4.3),
            SLIDE_WIDTH - MARGINS["left"] - MARGINS["right"], Inches(1.8)
        )
        company_frame = company_box.text_frame
        company_frame.word_wrap = True
        p = company_frame.paragraphs[0]
        p.text = COMPANY_NAME
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.font.name = FONTS["title"]
        p.alignment = PP_ALIGN.CENTER

        address_parts = [part.strip() for part in COMPANY_ADDRESS.split(",") if part.strip()]
        address_line_1 = ", ".join(address_parts[:3]) if address_parts else COMPANY_ADDRESS
        address_line_2 = ", ".join(address_parts[3:]) if len(address_parts) > 3 else ""

        p_address = company_frame.add_paragraph()
        p_address.text = address_line_1
        p_address.font.size = FONT_SIZES["body"]
        p_address.font.color.rgb = COLORS["text"]
        p_address.font.name = FONTS["body"]
        p_address.alignment = PP_ALIGN.CENTER

        if address_line_2:
            p_address_2 = company_frame.add_paragraph()
            p_address_2.text = address_line_2
            p_address_2.font.size = FONT_SIZES["body"]
            p_address_2.font.color.rgb = COLORS["text"]
            p_address_2.font.name = FONTS["body"]
            p_address_2.alignment = PP_ALIGN.CENTER

        p_website = company_frame.add_paragraph()
        p_website.text = COMPANY_WEBSITE
        p_website.font.size = FONT_SIZES["body"]
        p_website.font.color.rgb = COLORS["accent"]
        p_website.font.name = FONTS["body"]
        p_website.alignment = PP_ALIGN.CENTER

        logger.info("Added thank you slide: %s", country_name)

    def save(self, file_path: str):
        """Save presentation to file."""
        try:
            self.prs.save(file_path)
            logger.info("Presentation saved to %s", file_path)
        except Exception as e:
            logger.error("Error saving presentation: %s", str(e))
            raise
